#!/usr/bin/env python3
"""
generar_ppt_cpnb.py — Bot de Telegram + Generador de PPT CPNB-ZULIA
====================================================================

Lo que hace:
  1) Se conecta a Telegram con tu BOT y revisa los DOS grupos
     configurados en config.json:
       - chat_id_gubernamentales  → 10 cuadros (slide 1)
       - chat_id_consulados       →  8 cuadros (slide 2)
  2) Cuando alguien sube una foto, lee el PIE DE FOTO (caption)
     y lo compara con los nombres de los lugares conocidos.
       · Tolera tildes, mayúsculas y abreviaturas
         ("PSUV", "psuv", "Palacio de Justicia",
          "palacio justicia", "Consulado España", "espanol", etc.)
       · Si llega otra foto para el mismo lugar y es MÁS RECIENTE,
         sustituye a la anterior.
       · Si no logra reconocer el caption, la deja en
         <grupo>/sin_clasificar/  y avisa al final.
  3) Genera CPNB_ZULIA.pptx con cada foto contenida y centrada
     en su cuadro (sin desbordes, conservando proporciones).

Uso:
    python generar_ppt_cpnb.py            # procesa Telegram y envía el reporte
    python generar_ppt_cpnb.py --ver-chats   # te dice los chat_id
    python generar_ppt_cpnb.py --generar-ppt # genera el PPT antiguo solo para compatibilidad

Dependencias (instalar UNA vez):
    pip install python-pptx requests Pillow
"""

import argparse
import difflib
import json
import os
import re
import sys
import unicodedata
import time
from pathlib import Path

# ---------------------------------------------------------------------------
def _faltante(nombre, paquete):
    print(f"ERROR: falta la librería '{nombre}'.")
    print(f"Instálala con:  pip install {paquete}")
    sys.exit(1)

try:
    import requests
except ImportError:
    _faltante("requests", "requests")
try:
    from PIL import Image
except ImportError:
    _faltante("Pillow", "Pillow")
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    _faltante("python-pptx", "python-pptx")


# ---------------------------------------------------------------------------
# Rutas y constantes
# ---------------------------------------------------------------------------
RAIZ = Path(__file__).parent.resolve()
# En Railway se define DATA_DIR=/data (volumen persistente); localmente usa RAIZ
_DATA = Path(os.environ.get("DATA_DIR", str(RAIZ)))
CONFIG_FILE  = RAIZ / "config.json"
ESTADO_FILE  = _DATA / "estado_telegram.json"
CARPETA_GUB  = _DATA / "fotos_gubernamentales"
CARPETA_CONS = _DATA / "fotos_consulados"
SALIDA_PPTX  = _DATA / "CPNB_ZULIA.pptx"
EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp"}


# ---------------------------------------------------------------------------
# Definición de los 18 lugares
#   slug       → identificador interno (nombre del archivo de foto)
#   nombre     → texto que va en el cuadro del PPT
#   coords     → texto con las coordenadas
#   keywords   → frases que se aceptan como "etiqueta" en el caption
#                (¡añade aquí cualquier alias que use tu equipo!)
# ---------------------------------------------------------------------------
LUGARES_GUB = [
    {"slug": "01_psuv",
     "nombre": "PSUV",
     "coords": "10.669693, -71.607350",
     "keywords": ["psuv"]},

    {"slug": "02_contraloria",
     "nombre": "Contraloría del Estado Zulia",
     "coords": "10.639110, -71.605094",
     "keywords": ["contraloria", "contraloria zulia", "contraloria estado",
                  "contraloria del estado zulia"]},

    {"slug": "03_alcaldia",
     "nombre": "Alcaldía del Municipio Maracaibo",
     "coords": "10.64111, -71.60758",
     "keywords": ["alcaldia", "alcaldia maracaibo",
                  "alcaldia del municipio maracaibo", "alcaldia municipio"]},

    {"slug": "04_bcv",
     "nombre": "BCV del Municipio Maracaibo",
     "coords": "10.64111, -71.60758",
     "keywords": ["bcv", "banco central", "banco central de venezuela",
                  "bcv maracaibo"]},

    {"slug": "05_gobernacion",
     "nombre": "Gobernación del Estado Zulia",
     "coords": "10.64226, -71.60758",
     "keywords": ["gobernacion", "gobernacion zulia",
                  "gobernacion del estado zulia"]},

    {"slug": "06_palacio_legislativo",
     "nombre": "Palacio Legislativo",
     "coords": "10.6421, -71.60723",
     "keywords": ["palacio legislativo", "legislativo", "consejo legislativo"]},

    {"slug": "07_palacio_justicia",
     "nombre": "Palacio de Justicia",
     "coords": "10.64307, -71.6176",
     "keywords": ["palacio de justicia", "palacio justicia", "tribunales",
                  "tsj"]},

    {"slug": "08_ministerio_publico",
     "nombre": "Ministerio Público",
     "coords": "10.66390, -71.61545",
     "keywords": ["ministerio publico", "ministerio publico zulia",
                  "fiscalia", "mp"]},

    {"slug": "09_seniat",
     "nombre": "Seniat",
     "coords": "10.66511, -71.61452",
     "keywords": ["seniat"]},

    {"slug": "10_pdvsa",
     "nombre": "PDVSA",
     "coords": "10.6423, -71.6151",
     "keywords": ["pdvsa"]},
]

LUGARES_CONS = [
    {"slug": "01_consulado_ecuador",
     "nombre": "Consulado Ecuatoriano",
     "coords": "10.67977, -71.60658",
     "keywords": ["consulado ecuatoriano", "consulado ecuador",
                  "ecuador", "ecuatoriano"]},

    {"slug": "02_consulado_alemania",
     "nombre": "Consulado de Alemania",
     "coords": "10.67453, -71.604846",
     "keywords": ["consulado alemania", "consulado de alemania",
                  "alemania", "aleman"]},

    {"slug": "03_consulado_colombia",
     "nombre": "Consulado de Colombia",
     "coords": "10.64111, -71.60758",
     "keywords": ["consulado colombia", "consulado de colombia",
                  "colombia", "colombiano"]},

    {"slug": "04_consulado_chile",
     "nombre": "Consulado de Chile",
     "coords": "10.69571, -71.59743",
     "keywords": ["consulado chile", "consulado de chile",
                  "chile", "chileno"]},

    {"slug": "05_consulado_espana",
     "nombre": "Consulado Español",
     "coords": "10.70785, -71.6089",
     "keywords": ["consulado espanol", "consulado espana",
                  "consulado de espana", "espana", "espanol", "espanha"]},

    {"slug": "06_consulado_guatemala",
     "nombre": "Consulado de Guatemala",
     "coords": "10.67102, -71.2351",
     "keywords": ["consulado guatemala", "consulado de guatemala",
                  "guatemala", "guatemalteco"]},

    {"slug": "07_consulado_italia",
     "nombre": "Consulado Italiano",
     "coords": "10.6674, -71.6249",
     "keywords": ["consulado italiano", "consulado italia",
                  "consulado de italia", "italia", "italiano"]},

    {"slug": "08_consulado_portugal",
     "nombre": "Consulado de Portugal",
     "coords": "10.668, -71.625",
     "keywords": ["consulado portugal", "consulado de portugal",
                  "portugal", "portugues"]},
]


# ---------------------------------------------------------------------------
# Normalización y matching de captions
# ---------------------------------------------------------------------------
def normalizar(texto: str) -> str:
    """Pasa a minúsculas, quita tildes y deja solo letras/números/espacios."""
    if not texto:
        return ""
    t = unicodedata.normalize("NFKD", texto)
    t = "".join(c for c in t if not unicodedata.combining(c))
    t = t.lower()
    t = re.sub(r"[^a-z0-9\s]", " ", t)
    t = re.sub(r"\s+", " ", t).strip()
    return t


def construir_indice():
    """Crea un mapa {keyword_normalizado: lugar} para búsquedas rápidas."""
    indice = {}
    for grupo, lista in (("gub", LUGARES_GUB), ("cons", LUGARES_CONS)):
        for lugar in lista:
            for kw in lugar["keywords"]:
                indice[normalizar(kw)] = (grupo, lugar)
    return indice


INDICE = construir_indice()
TODAS_LAS_KEYS = list(INDICE.keys())


def matchear_caption(caption: str, grupo_chat: str):
    """Devuelve (grupo, lugar) si reconoce el caption; None si no.
       grupo_chat = 'gub' o 'cons' (de qué grupo viene la foto)."""
    n = normalizar(caption)
    if not n:
        return None

    # 1) Match exacto por substring (preferimos la keyword más larga)
    candidatos = []
    for kw, (grupo, lugar) in INDICE.items():
        if kw in n or n in kw:
            candidatos.append((len(kw), grupo, lugar))
    if candidatos:
        # Si hay candidatos del MISMO grupo del chat, los preferimos
        del_grupo = [c for c in candidatos if c[1] == grupo_chat]
        elegidos = del_grupo if del_grupo else candidatos
        elegidos.sort(reverse=True)              # keyword más larga gana
        return (elegidos[0][1], elegidos[0][2])

    # 2) Fuzzy match (tolera erratas leves)
    cercanos = difflib.get_close_matches(n, TODAS_LAS_KEYS, n=1, cutoff=0.78)
    if cercanos:
        return INDICE[cercanos[0]]

    return None


# ---------------------------------------------------------------------------
# Configuración / estado
# ---------------------------------------------------------------------------
def cargar_config():
    if not CONFIG_FILE.is_file():
        print("Falta el archivo 'config.json'. Crea uno con:\n")
        print(json.dumps({
            "bot_token": "PEGA_AQUI_EL_TOKEN_DEL_BOT",
            "chat_id_gubernamentales": 0,
            "chat_id_consulados": 0,
            "chat_id_reporte": 0,
        }, indent=2, ensure_ascii=False))
        sys.exit(1)
    with CONFIG_FILE.open(encoding="utf-8") as f:
        return json.load(f)


def cargar_estado():
    if ESTADO_FILE.is_file():
        with ESTADO_FILE.open(encoding="utf-8") as f:
            return json.load(f)
    return {"last_update_id": 0, "fotos": {}}   # fotos: slug -> {ts, msg_id}


def guardar_estado(estado):
    with ESTADO_FILE.open("w", encoding="utf-8") as f:
        json.dump(estado, f, indent=2, ensure_ascii=False)


# ---------------------------------------------------------------------------
# Telegram
# ---------------------------------------------------------------------------
def tg_get_updates(token, offset=None, timeout=15):
    url = f"https://api.telegram.org/bot{token}/getUpdates"
    params = {"timeout": timeout, "allowed_updates": '["message"]'}
    if offset is not None:
        params["offset"] = offset
    r = requests.get(url, params=params, timeout=timeout + 10)
    r.raise_for_status()
    data = r.json()
    if not data.get("ok"):
        raise RuntimeError(f"Telegram error: {data}")
    return data["result"]


def tg_get_file_path(token, file_id):
    r = requests.get(f"https://api.telegram.org/bot{token}/getFile",
                     params={"file_id": file_id}, timeout=30)
    r.raise_for_status()
    return r.json()["result"]["file_path"]


def tg_descargar_archivo(token, file_path, destino: Path):
    url = f"https://api.telegram.org/file/bot{token}/{file_path}"
    r = requests.get(url, timeout=60)
    r.raise_for_status()
    destino.write_bytes(r.content)


def tg_send_message(token, chat_id, text):
    url = f"https://api.telegram.org/bot{token}/sendMessage"
    r = requests.post(url, data={
        "chat_id": chat_id,
        "text": text,
        "parse_mode": "HTML",
    }, timeout=30)
    r.raise_for_status()


def tg_send_photo(token, chat_id, photo_path, caption=None):
    url = f"https://api.telegram.org/bot{token}/sendPhoto"
    with open(photo_path, "rb") as f:
        files = {"photo": f}
        data = {"chat_id": chat_id}
        if caption:
            data["caption"] = caption
        r = requests.post(url, files=files, data=data, timeout=60)
    r.raise_for_status()


def tg_send_media_group(token, chat_id, photo_paths, caption=None):
    url = f"https://api.telegram.org/bot{token}/sendMediaGroup"
    for i in range(0, len(photo_paths), 10):
        chunk = photo_paths[i:i+10]
        files = {}
        media = []
        for idx, path in enumerate(chunk):
            key = f"file{idx}"
            item = {"type": "photo", "media": f"attach://{key}"}
            if idx == 0 and caption:
                item["caption"] = caption
                item["parse_mode"] = "HTML"
            media.append(item)
            files[key] = open(path, "rb")
        try:
            data = {
                "chat_id": chat_id,
                "media": json.dumps(media),
            }
            r = requests.post(url, data=data, files=files, timeout=60)
            r.raise_for_status()
        finally:
            for f in files.values():
                f.close()


def es_reporte_valido(grupo_match, slug):
    # Considerar reporte válido para consulados o para el grupo gubernamentales
    # (antes solo se disparaba para el slug 08_ministerio_publico)
    return grupo_match == "cons" or grupo_match == "gub"


def carpeta_de(grupo):
    return CARPETA_GUB if grupo == "gub" else CARPETA_CONS


def borrar_fotos_anteriores(carpeta: Path, slug: str):
    """Elimina cualquier foto previa de ese lugar (para sustituir)."""
    if not carpeta.is_dir():
        return
    for p in carpeta.iterdir():
        if p.is_file() and p.stem == slug:
            try:
                p.unlink()
            except Exception:
                pass


def procesar_telegram(config, estado):
    token = config["bot_token"]
    cid_gub = config.get("chat_id_gubernamentales")
    cid_cons = config.get("chat_id_consulados")
    reporte_chat_id = config.get("chat_id_reporte")

    if not token or token.startswith("PEGA_"):
        print("config.json: pega el bot_token primero. Saltando Telegram.")
        return [], []
    if not cid_gub or not cid_cons:
        print("config.json: faltan chat_id_gubernamentales / chat_id_consulados.")
        print("Ejecuta primero:  python generar_ppt_cpnb.py --ver-chats")
        return [], []

    CARPETA_GUB.mkdir(exist_ok=True)
    CARPETA_CONS.mkdir(exist_ok=True)
    (CARPETA_GUB / "sin_clasificar").mkdir(exist_ok=True)
    (CARPETA_CONS / "sin_clasificar").mkdir(exist_ok=True)

    offset = estado.get("last_update_id", 0) + 1
    print(f"Consultando Telegram (offset={offset})...")
    updates = tg_get_updates(token, offset=offset, timeout=10)
    print(f"  {len(updates)} actualizaciones recibidas.")

    sin_clasif = []
    sustituciones = []
    reporte_trigger = []
    reporte_fotos_cons = []
    reporte_fotos_min = []

    for upd in updates:
        estado["last_update_id"] = max(estado["last_update_id"], upd["update_id"])
        msg = upd.get("message") or upd.get("channel_post")
        if not msg:
            continue

        chat_id = msg["chat"]["id"]
        if chat_id == cid_gub:
            grupo_chat = "gub"
        elif chat_id == cid_cons:
            grupo_chat = "cons"
        else:
            continue

        photos = msg.get("photo")
        if not photos:
            continue

        caption = msg.get("caption", "") or ""
        ts = msg["date"]
        mid = msg["message_id"]
        biggest = max(photos, key=lambda p: p.get("file_size", 0))
        file_id = biggest["file_id"]
        carpeta = carpeta_de(grupo_chat)

        try:
            fp = tg_get_file_path(token, file_id)
            ext = Path(fp).suffix.lower() or ".jpg"
        except Exception as e:
            print(f"  ! Error obteniendo file_path: {e}")
            continue

        match = matchear_caption(caption, grupo_chat)

        if match is None:
            # Sin etiqueta o no reconocida
            destino = carpeta / "sin_clasificar" / f"{ts:010d}_{mid:09d}{ext}"
            try:
                if not destino.exists():
                    tg_descargar_archivo(token, fp, destino)
                sin_clasif.append((grupo_chat, caption, destino.name))
                print(f"  ? sin_clasificar/{destino.name}  (caption: {caption!r})")
            except Exception as e:
                print(f"  ! Error descargando: {e}")
            continue

        grupo_match, lugar = match
        # Si el caption señala otro grupo (ej: alguien manda PSUV al
        # grupo de consulados), respetamos el caption pero lo enviamos
        # a la carpeta correcta.
        carpeta_destino = carpeta_de(grupo_match)
        slug = lugar["slug"]
        previo = estado["fotos"].get(slug)

        if previo and previo["ts"] >= ts:
            # Ya tenemos una foto MÁS reciente o igual para ese lugar
            print(f"  · {lugar['nombre']}: ignorada (más antigua que la actual)")
            continue

        try:
            destino = carpeta_destino / f"{slug}{ext}"
            borrar_fotos_anteriores(carpeta_destino, slug)
            tg_descargar_archivo(token, fp, destino)
            estado["fotos"][slug] = {"ts": ts, "msg_id": mid,
                                     "caption": caption, "grupo": grupo_match}
            if previo:
                sustituciones.append(lugar["nombre"])

            tipo = "Consulado" if grupo_match == "cons" else "Ministerio Público" if slug == "08_ministerio_publico" else "Otro"
            print(f"  ✓ {lugar['nombre']} ({tipo})  →  {destino.name}")

            if es_reporte_valido(grupo_match, slug):
                tipo_reporte = "cons" if grupo_match == "cons" else "min"
                reporte_trigger.append(tipo_reporte)
                reporte_trigger.append(tipo_reporte)
                # No asignar `reporte_chat_id` desde el chat origen.
                # Los reportes se envían únicamente al `chat_id_reporte`
                # definido en `config.json`.
                if not reporte_chat_id:
                    reporte_chat_id = chat_id
        except Exception as e:
            print(f"  ! Error guardando {slug}: {e}")

    guardar_estado(estado)

    if sustituciones:
        print(f"\nSe actualizaron {len(sustituciones)} fotos: "
              + ", ".join(sustituciones))

    if sin_clasif:
        print(f"\n⚠  {len(sin_clasif)} fotos sin clasificar:")
        for grupo, cap, nombre in sin_clasif:
            cg = "gubernamentales" if grupo == "gub" else "consulados"
            print(f"    - {cg}/sin_clasificar/{nombre}   caption: {cap!r}")
        print("   Renómbralas manualmente (con el slug del lugar) y vuelve")
        print("   a ejecutar:  python generar_ppt_cpnb.py --generar-ppt")

    if reporte_trigger:
        if not reporte_chat_id:
            print("Aviso: no se definió chat_id_reporte en config.json; no se puede enviar el reporte.")
        else:
            tipos = set(reporte_trigger)
            r1, r2 = elegir_dos_responsables()

            if "cons" in tipos:
                try:
                    fotos_cons = []
                    for lugar in LUGARES_CONS:
                        p = buscar_foto_de(lugar["slug"], CARPETA_CONS)
                        if p:
                            fotos_cons.append(str(p))
                    texto = texto_reporte_consulados(r1)
                    if fotos_cons:
                        tg_send_media_group(token, reporte_chat_id, fotos_cons, caption=texto)
                    else:
                        tg_send_message(token, reporte_chat_id, texto)
                    print(f"\nReporte de consulados enviado a chat_id={reporte_chat_id}.")
                except Exception as e:
                    print(f"\n! Error enviando reporte de consulados a Telegram: {e}")

            if "min" in tipos:
                try:
                    fotos_min = []
                    for lugar in LUGARES_GUB:
                        p = buscar_foto_de(lugar["slug"], CARPETA_GUB)
                        if p:
                            fotos_min.append(str(p))
                    texto = texto_reporte_ministerio(r2)
                    if fotos_min:
                        tg_send_media_group(token, reporte_chat_id, fotos_min, caption=texto)
                    else:
                        tg_send_message(token, reporte_chat_id, texto)
                    print(f"\nReporte de Ministerio Público enviado a chat_id={reporte_chat_id}.")
                except Exception as e:
                    print(f"\n! Error enviando reporte de Ministerio Público a Telegram: {e}")

    return sin_clasif, sustituciones


def ver_chats(config):
    token = config["bot_token"]
    if not token or token.startswith("PEGA_"):
        print("Pega antes el bot_token en config.json.")
        return
    print("Pidiendo updates a Telegram...")
    print("(Si no aparece nada, escribe cualquier mensaje en cada grupo y reintenta.)\n")
    updates = tg_get_updates(token, timeout=10)
    vistos = {}
    for upd in updates:
        msg = upd.get("message") or upd.get("channel_post")
        if not msg:
            continue
        chat = msg["chat"]
        vistos[chat["id"]] = chat.get("title") or chat.get("username") or "(sin título)"
    if not vistos:
        print("No se vieron grupos. Escribe un mensaje en cada grupo y reintenta.")
        return
    print("Grupos detectados:")
    for cid, nombre in vistos.items():
        print(f"  chat_id = {cid:<15}   →  {nombre}")


# ---------------------------------------------------------------------------
# Layout del PPT
# ---------------------------------------------------------------------------
SLIDE_W = 13.333
SLIDE_H = 7.5
COLS, ROWS = 5, 2
MARGIN_X, MARGIN_Y = 0.30, 0.30
GAP_X, GAP_Y = 0.12, 0.18

CARD_W = (SLIDE_W - 2 * MARGIN_X - (COLS - 1) * GAP_X) / COLS
CARD_H = (SLIDE_H - 2 * MARGIN_Y - (ROWS - 1) * GAP_Y) / ROWS

PAD = 0.08
HEADER_H = 0.32
LABEL_H = 0.50
COORD_H = 0.36
PHOTO_H = CARD_H - (PAD * 5 + HEADER_H + LABEL_H + COORD_H)

NEGRO = RGBColor(0, 0, 0)
BLANCO = RGBColor(255, 255, 255)

# Lista de responsables — se rota automáticamente por día
PERSONAS = [
    {"nombre": "Adrian Eduardo", "apellido": "Soto Portillo", "ci": "20660371", "telefono": "04129623811"},
    {"nombre": "Orianys Chiquinquira", "apellido": "Fernández Delgado", "ci": "34532279", "telefono": "04147124410"},
    {"nombre": "Royer junior", "apellido": "Rodríguez Royer", "ci": "30999326", "telefono": "04146106701"},
    {"nombre": "Robersy Grismaily", "apellido": "Millano López", "ci": "32514356", "telefono": "04129516659"},
    {"nombre": "Yoswar José", "apellido": "Guerrero Viloria", "ci": "30355447", "telefono": "0424-6232296"},
    {"nombre": "Osman José", "apellido": "Fernández Carrillo", "ci": "30149787", "telefono": "04246440658"},
    {"nombre": "Jehan Robert", "apellido": "Matos Acosta", "ci": "27395099", "telefono": "04123007653"},
    {"nombre": "Jhon Albert", "apellido": "Gerardino García", "ci": "30182381", "telefono": "04121768042"},
    {"nombre": "Anadelis Beatriz", "apellido": "Fernández Carrillo", "ci": "30149754", "telefono": "04246871865"},
]


def elegir_responsable(fecha=None):
    """Devuelve el primer responsable del día (compatibilidad)."""
    r1, _ = elegir_dos_responsables(fecha)
    return r1


def elegir_dos_responsables(fecha=None):
    """Devuelve (r1, r2) — dos responsables del día, rotan por día."""
    from datetime import date
    if fecha is None:
        fecha = date.today()
    n = len(PERSONAS)
    idx1 = fecha.toordinal() % n
    idx2 = (idx1 + 1) % n
    return PERSONAS[idx1], PERSONAS[idx2]


def elegir_responsable_ejecucion():
    """Selecciona un responsable distinto por ejecución (basado en time_ns)."""
    t = time.time_ns()
    idx = int(t % len(PERSONAS))
    return PERSONAS[idx]

def texto_reporte_consulados(r):
    return (
        "CPNB-DAET-DIE- ZULIA\n\n"
        "📌Ubicación: Estado Zulia, Municipio Maracaibo, Parroquia Coquivacoa, "
        "Raul Leoni, Idelfonso Vasquez, Olegario Villalobos.\n\n"
        "📌Situación: Se le informa a la Superioridad, sobre el monitoreo estratégico, "
        "en los siguientes Consulados: Italia, Guatemala, Portugal, Alemania, Chile, "
        "Ecuador, Colombia, España. Para la hora sin alteraciones de Orden.\n\n"
        "RESPONSABLE DEL RECORRIDO:\n"
        "JERARQUÍA: OFICIAL\n"
        f"NOMBRE: {r['nombre']}\n"
        f"APELLIDO: {r['apellido']}\n"
        f"CÉDULA: {r['ci']}\n"
        f"TLF: {r['telefono']}\n"
        "SEGUIREMOS INFORMANDO"
    )


def texto_reporte_ministerio(r):
    return (
        "CPNB-DAET-DIE- ZULIA\n\n"
        "📌Ubicación: Estado Zulia, Municipio Maracaibo, Parroquia Coquivacoa, "
        "Raul Leoni, Idelfonso Vasquez, Olegario Villalobos.\n\n"
        "📌Situación: Se le informa a la Superioridad, sobre el monitoreo estratégico, "
        "en los diferentes entes públicos (Defensa Pública, Alcaldía, SENIAT, Contraloría, "
        "Palacio de Justicia, Banco Central de Venezuela, PDVSA, PSUV). "
        "Para la hora sin alteraciones de Orden.\n\n"
        "RESPONSABLE DEL RECORRIDO:\n"
        "JERARQUÍA: OFICIAL\n"
        f"NOMBRE: {r['nombre']}\n"
        f"APELLIDO: {r['apellido']}\n"
        f"CÉDULA: {r['ci']}\n"
        f"TLF: {r['telefono']}\n"
        "SEGUIREMOS INFORMANDO"
    )


def caja_texto(slide, x, y, w, h, texto, font_size=10, bold=False,
               fill=BLANCO, border=True):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if border:
        box.line.color.rgb = NEGRO
        box.line.width = Pt(0.75)
    else:
        box.line.fill.background()
    box.shadow.inherit = False

    tf = box.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = NEGRO


def insertar_foto_contenida(slide, foto_path, box_x, box_y, box_w, box_h):
    try:
        with Image.open(foto_path) as img:
            iw, ih = img.size
    except Exception as e:
        print(f"  ! No se pudo leer {foto_path}: {e}")
        return
    if iw == 0 or ih == 0:
        return

    box_ratio = box_w / box_h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        new_w = box_w
        new_h = box_w / img_ratio
    else:
        new_h = box_h
        new_w = box_h * img_ratio
    off_x = box_x + (box_w - new_w) / 2
    off_y = box_y + (box_h - new_h) / 2
    slide.shapes.add_picture(str(foto_path),
                             Inches(off_x), Inches(off_y),
                             width=Inches(new_w), height=Inches(new_h))


def dibujar_cuadro(slide, col, row, nombre, coords, foto_path):
    x = MARGIN_X + col * (CARD_W + GAP_X)
    y = MARGIN_Y + row * (CARD_H + GAP_Y)

    outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y),
                                   Inches(CARD_W), Inches(CARD_H))
    outer.fill.background()
    outer.line.color.rgb = NEGRO
    outer.line.width = Pt(1.0)
    outer.shadow.inherit = False

    cy = y + PAD
    caja_texto(slide, x + PAD, cy, CARD_W - 2 * PAD, HEADER_H,
               "CPNB-ZULIA", font_size=11, bold=True)
    cy += HEADER_H + PAD

    if foto_path and Path(foto_path).is_file():
        insertar_foto_contenida(slide, foto_path,
                                x + PAD, cy,
                                CARD_W - 2 * PAD, PHOTO_H)
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x + PAD), Inches(cy),
                                    Inches(CARD_W - 2 * PAD), Inches(PHOTO_H))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(245, 245, 245)
        ph.line.color.rgb = RGBColor(200, 200, 200)
        ph.line.width = Pt(0.5)
        ph.shadow.inherit = False
    cy += PHOTO_H + PAD

    fs_nombre = 10 if len(nombre) <= 28 else 9 if len(nombre) <= 36 else 8
    caja_texto(slide, x + PAD, cy, CARD_W - 2 * PAD, LABEL_H,
               nombre, font_size=fs_nombre, bold=False)
    cy += LABEL_H + PAD

    caja_texto(slide, x + PAD, cy, CARD_W - 2 * PAD, COORD_H,
               coords, font_size=9, bold=False)


def buscar_foto_de(slug: str, carpeta: Path):
    """Devuelve la ruta de la foto guardada para ese slug, o None."""
    if not carpeta.is_dir():
        return None
    for p in carpeta.iterdir():
        if p.is_file() and p.stem == slug and p.suffix.lower() in EXTS:
            return p
    return None


def generar_ppt():
    pres = Presentation()
    pres.slide_width = Inches(SLIDE_W)
    pres.slide_height = Inches(SLIDE_H)
    blank = pres.slide_layouts[6]

    # Slide 1: gubernamentales (5x2)
    s1 = pres.slides.add_slide(blank)
    for i, lugar in enumerate(LUGARES_GUB):
        col = i % COLS
        row = i // COLS
        foto = buscar_foto_de(lugar["slug"], CARPETA_GUB)
        dibujar_cuadro(s1, col, row, lugar["nombre"], lugar["coords"], foto)

    # Slide 2: consulados (8 cuadros, mismo grid)
    s2 = pres.slides.add_slide(blank)
    for i, lugar in enumerate(LUGARES_CONS):
        col = i % COLS
        row = i // COLS
        foto = buscar_foto_de(lugar["slug"], CARPETA_CONS)
        dibujar_cuadro(s2, col, row, lugar["nombre"], lugar["coords"], foto)

    pres.save(SALIDA_PPTX)

    g_listas = sum(1 for l in LUGARES_GUB if buscar_foto_de(l["slug"], CARPETA_GUB))
    c_listas = sum(1 for l in LUGARES_CONS if buscar_foto_de(l["slug"], CARPETA_CONS))
    print(f"\nPPT listo → {SALIDA_PPTX}")
    print(f"  Gubernamentales: {g_listas}/{len(LUGARES_GUB)} cuadros con foto")
    print(f"  Consulados:      {c_listas}/{len(LUGARES_CONS)} cuadros con foto")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(description="Bot Telegram → reporte CPNB-DAET-DIE-ZULIA")
    parser.add_argument("--ver-chats", action="store_true",
                        help="Muestra los chat_id de los grupos donde está el bot")
    parser.add_argument("--generar-ppt", action="store_true",
                        help="Genera el PPT antiguo solo para compatibilidad")
    args = parser.parse_args()

    config = cargar_config()

    if args.ver_chats:
        ver_chats(config)
        return

    if args.generar_ppt:
        print("La opción --generar-ppt está deshabilitada: no se generarán archivos PPT.")
        print("El flujo actual envía únicamente la minuta con las fotos por Telegram.")
        return

    estado = cargar_estado()
    try:
        procesar_telegram(config, estado)
    except Exception as e:
        print(f"AVISO: fallo al consultar Telegram ({e}).")
        print("No se puede enviar el reporte sin conexión a Telegram.")


if __name__ == "__main__":
    main()